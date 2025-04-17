import os
import streamlit as st
import openai
from dotenv import load_dotenv
from fpdf import FPDF
from gtts import gTTS
from pptx import Presentation
import base64

# Subida a Google Drive
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive

# --- Configuración ---
load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")

st.title("🎥 Generador de Guiones VSL con IA")
st.write("Completa los datos para generar tu VSL en texto, PDF, audio y presentación.")

# --- Autenticación Google Drive (solo la primera vez) ---
@st.cache_resource
def autenticar_drive():
    gauth = GoogleAuth()
    gauth.LocalWebserverAuth()  # Abre navegador para login
    return GoogleDrive(gauth)

drive = autenticar_drive()

# --- Inputs ---
nicho = st.selectbox("Selecciona tu tipo de negocio:", [
    "Coach", "Terapeuta", "Consultor", "Formador", "Nutricionista",
    "Agencia", "Psicólogo", "Negocio SaaS", "Entrenador personal", "Otro"
])

nombre_producto   = st.text_input("Nombre del producto/servicio", key="nombre_producto")
publico_objetivo  = st.text_input("Público objetivo (Ej: coaches, madres…)", key="publico_objetivo")
dolor_problema    = st.text_area("Dolor o problema que resuelve", key="dolor_problema")
beneficios_clave  = st.text_area("Beneficios clave (separa con comas)", key="beneficios_clave")
precio_forma_pago = st.text_input("Precio / forma de pago", key="precio")
garantia          = st.text_input("Garantía ofrecida", key="garantia")
cta               = st.text_input("Llamada a la acción (CTA)", key="cta")

ejemplos_nicho = {
    "Coach": "Ejemplo: 'Con mi programa de coaching, mis clientes superan bloqueos y mejoran su vida profesional y personal...'",
    "Terapeuta": "Ejemplo: 'Mis pacientes reducen el estrés y la ansiedad, y recuperan su equilibrio emocional...'",
    "Consultor": "Ejemplo: 'Ayudo a negocios como el tuyo a aumentar su rentabilidad y tomar decisiones estratégicas con confianza...'",
    "Formador": "Ejemplo: 'Mis formaciones prácticas permiten aplicar lo aprendido desde el primer día, aumentando el rendimiento...'",
    "Nutricionista": "Ejemplo: 'Con mis planes personalizados, los clientes pierden peso sin pasar hambre y mejoran su salud...'",
    "Agencia": "Ejemplo: 'Nuestros clientes multiplican sus ventas con campañas bien diseñadas y estrategias de conversión...'",
    "Psicólogo": "Ejemplo: 'Acompaño a mis pacientes a resolver traumas, mejorar su autoestima y construir relaciones sanas...'",
    "Negocio SaaS": "Ejemplo: 'Nuestra plataforma ahorra tiempo, automatiza procesos y mejora la experiencia del cliente...'",
    "Entrenador personal": "Ejemplo: 'Mis clientes transforman su físico, recuperan energía y autoestima con rutinas personalizadas...'",
    "Otro": ""
}

# --- Botón de generación ---
if st.button("🎬 Generar Guion de VSL"):
    with st.spinner("Generando guion…"):
        prompt = f"""
Eres un copywriter profesional, experto en Video Sales Letters.
Genera un guion de VSL en ESPAÑOL de entre 600 y 800 palabras (máx. 7 minutos), con tono profesional y lenguaje sencillo.

Nicho: {nicho}
{ejemplos_nicho[nicho]}

Estructura obligatoria:
1. Hook inicial impactante.
2. Historia personal adaptada al nicho.
3. Dolor/problema específico.
4. Presentación del producto como solución.
5. Beneficios clave (bullet points).
6. Testimonios con nombre y profesión.
7. Objeciones comunes + rebatidas.
8. Garantía.
9. Llamada a la acción clara.

Datos personalizados:
- Producto: {nombre_producto}
- Público objetivo: {publico_objetivo}
- Problema principal: {dolor_problema}
- Beneficios: {beneficios_clave}
- Precio / forma de pago: {precio_forma_pago}
- Garantía: {garantia}
- Llamada a la acción: {cta}
"""

        try:
            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "Eres un experto copywriter en VSL."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=900,
                temperature=0.7
            )
            guion = response.choices[0].message.content

            st.subheader("📄 Guion generado")
            st.markdown(guion)

            # --- Descarga TXT local + subida a Drive ---
            txt_filename = "guion_vsl.txt"
            with open(txt_filename, "w", encoding="utf-8") as f:
                f.write(guion)

            st.download_button("📄 Descargar Guion en TXT", data=guion, file_name=txt_filename, mime="text/plain")

            # Subir a Drive
            archivo_drive = drive.CreateFile({'title': txt_filename})
            archivo_drive.SetContentFile(txt_filename)
            archivo_drive.Upload()
            enlace_drive = f"https://drive.google.com/file/d/{archivo_drive['id']}/view"
            st.success("✅ Guion subido a Google Drive correctamente.")
            st.markdown(f"🔗 [Ver en Google Drive]({enlace_drive})")

            # --- PDF ---
            def crear_pdf(texto):
                pdf = FPDF()
                pdf.add_page()
                pdf.set_font("Arial", size=12)
                for linea in texto.split('\n'):
                    pdf.multi_cell(0, 10, txt=linea.encode('latin-1', 'replace').decode('latin-1'))
                return pdf.output(dest='S').encode('latin-1')

            pdf_data = crear_pdf(guion)
            b64 = base64.b64encode(pdf_data).decode('latin-1')
            href = f'<a href="data:application/octet-stream;base64,{b64}" download="guion_vsl.pdf">📥 Descargar Guion en PDF</a>'
            st.markdown(href, unsafe_allow_html=True)

            # --- Audio narrado ---
            try:
                tts = gTTS(text=guion, lang='es')
                tts.save("guion_vsl.mp3")
                with open("guion_vsl.mp3", "rb") as audio_file:
                    audio_bytes = audio_file.read()

                st.subheader("🎙️ Escucha tu guion narrado")
                st.audio(audio_bytes, format="audio/mp3")
                st.download_button("🎧 Descargar narración MP3", data=audio_bytes, file_name="guion_vsl.mp3", mime="audio/mpeg")

            except Exception as e:
                st.warning("No se pudo generar la narración.")
                st.text(str(e))

            # --- PowerPoint ---
            try:
                prs = Presentation()
                layout = prs.slide_layouts[1]
                bloques = guion.split('\n\n')
                for i, bloque in enumerate(bloques):
                    if bloque.strip() == "":
                        continue
                    slide = prs.slides.add_slide(layout)
                    slide.shapes.title.text = f"Slide {i+1}"
                    texto_slide = bloque.strip().replace('\n', ' ')
                    if len(texto_slide) > 250:
                        texto_slide = texto_slide[:247] + '...'
                    slide.placeholders[1].text = texto_slide

                prs.save("guion_vsl.pptx")
                with open("guion_vsl.pptx", "rb") as pptx_file:
                    pptx_bytes = pptx_file.read()

                st.subheader("📊 Presentación generada")
                st.download_button("📥 Descargar presentación PPTX", data=pptx_bytes, file_name="guion_vsl.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

            except Exception as e:
                st.warning("No se pudo generar la presentación.")
                st.text(str(e))

        except Exception as e:
            st.error(f"Ocurrió un error al generar el guion: {e}")
